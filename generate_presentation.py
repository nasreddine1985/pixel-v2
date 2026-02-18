#!/usr/bin/env python3
"""
PIXEL-V2 PowerPoint Presentation Generator
Generates a professional PowerPoint presentation for the PIXEL-V2 project
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_shape.text_frame.paragraphs[0].font.size = Pt(54)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

def create_section_slide(prs, title):
    """Create a section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(48)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def create_content_slide(prs, title, content_list):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    text_frame = content_shape.text_frame
    text_frame.clear()
    
    for item in content_list:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
            
        p = text_frame.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = Pt(20 - (level * 2))

def create_two_column_slide(prs, title, left_content, right_content):
    """Create a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.5), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    for item in left_content:
        p = left_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.5), Inches(4.5), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    for item in right_content:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.space_after = Pt(10)

def main():
    """Generate the PIXEL-V2 PowerPoint presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    create_title_slide(prs, 
                      "PIXEL-V2",
                      "Enterprise Payment Message Processing System")
    
    # Slide 2: Executive Summary
    create_content_slide(prs,
                        "Executive Summary",
                        [
                            "Enterprise-grade multi-flow payment processing platform",
                            "Processing Capability: Multiple payment flows & schemes",
                            "Architecture: Microservices with Apache Camel",
                            "Technology: Spring Boot 3.4.1, Java 21, Camel 4.1.0",
                            "Infrastructure: Containerized with Docker",
                            "Integration: Kafka, IBM MQ, Redis, PostgreSQL, CFT",
                            "Compliance: ISO 20022 standards compliant"
                        ])
    
    # Slide 3: What is PIXEL-V2?
    create_content_slide(prs,
                        "What is PIXEL-V2?",
                        [
                            "Multi-flow payment message processing platform:",
                            ("Supports multiple payment flows and schemes", 1),
                            ("Receives payment messages from multiple sources (JMS)", 1),
                            ("Validates messages against XSD schemas (ISO 20022)", 1),
                            ("Transforms messages using XSLT templates", 1),
                            ("Enriches data with referential information", 1),
                            ("Archives messages to NAS file systems", 1),
                            ("Logs processing events to Kafka for monitoring", 1)
                        ])
    
    # Slide 4: Business Value
    create_content_slide(prs,
                        "Business Value",
                        [
                            "🚀 High Performance",
                            ("Spring caching for sub-millisecond referential lookups", 1),
                            "🔒 Reliability",
                            ("Comprehensive validation and error handling", 1),
                            "📊 Observability",
                            ("Complete audit trail via Kafka event streaming", 1),
                            "🔄 Flexibility",
                            ("Configurable flows and transformation rules", 1),
                            "💾 Compliance",
                            ("Automatic message archiving for regulatory requirements", 1)
                        ])
    
    # Slide 5: System Architecture
    create_section_slide(prs, "System Architecture")
    
    # Slide 6: Multi-Module Structure
    create_content_slide(prs,
                        "Multi-Module Maven Project",
                        [
                            "PIXEL-V2 (Parent Module)",
                            ("technical-framework/ - Reusable Camel Kamelets", 1),
                            ("referential/ - Configuration Service API", 1),
                            ("flow-ch/ - ICHSIC Flow (Example Implementation)", 1),
                            ("flow-xxx/ - Future flows (SEPA, SWIFT, etc.)", 1),
                            "",
                            "Generic Processing Pipeline:",
                            ("JMS Queue → k-mq-starter → k-identification", 1),
                            ("→ k-xsd-validation → k-xsl-transformation", 1),
                            ("→ k-log-flow-summary", 1)
                        ])
    
    # Slide 7: Key Components
    create_section_slide(prs, "Key Components")
    
    # Slide 8: Technical Framework - Kamelets
    create_content_slide(prs,
                        "Technical Framework - Custom Kamelets",
                        [
                            "k-mq-starter: JMS queue listener & flow ID generation",
                            "k-identification: Payment ID & referential lookup (Spring Cache)",
                            "k-xsd-validation: XML Schema validation (ISO 20022)",
                            "k-xsl-transformation: XSLT-based message transformation",
                            "k-log-flow-summary: Audit logging to Kafka",
                            "k-kafka-publisher: Asynchronous event streaming",
                            "k-cft-publisher: CFT file transfer integration",
                            "k-http-publisher: RESTful service integration",
                            "k-db-tx: Database transaction management",
                            "k-duplicate-check: Duplicate message detection"
                        ])
    
    # Slide 9: Flow-CH Component
    create_content_slide(prs,
                        "Flow-CH: ICHSIC Flow Example",
                        [
                            "One of many payment flows in PIXEL-V2",
                            "Flow Code: ICHSIC (Switzerland IN SEPA SIC)",
                            "",
                            "Input: JMS queues (IBM MQ)",
                            "Validation: pacs.008.001.02.ch.02.xsd",
                            "Transformation: overall-xslt-ch-pacs008-001-02.xsl",
                            "Output: Kafka topics, NAS archive, CFT transfer",
                            "Configuration: Dynamic via referential service",
                            "",
                            "Each flow module follows the same pattern",
                            "Reusable Kamelets from technical-framework"
                        ])
    
    # Slide 10: Referential Service
    create_content_slide(prs,
                        "Referential Configuration Service",
                        [
                            "Central configuration repository for ALL flows",
                            "",
                            "Features:",
                            ("Flow definitions (ICHSIC, SEPA, SWIFT, etc.)", 1),
                            ("Country and partner mappings", 1),
                            ("Business rules and validation parameters", 1),
                            ("Character encoding configurations", 1),
                            ("Real-time configuration retrieval", 1),
                            "",
                            "Technology:",
                            ("Spring Boot 3.4.1 + PostgreSQL + Liquibase", 1)
                        ])
    
    # Slide 11: Technical Stack
    create_section_slide(prs, "Technical Stack")
    
    # Slide 11.5: Multi-Flow Architecture Design
    create_content_slide(prs,
                        "Multi-Flow Architecture Design",
                        [
                            "Scalable architecture supporting multiple payment flows:",
                            "",
                            "Shared Components:",
                            ("✓ Technical Framework (reusable Kamelets)", 1),
                            ("✓ Referential Service (centralized configuration)", 1),
                            ("✓ Infrastructure (Kafka, Redis, PostgreSQL, etc.)", 1),
                            "",
                            "Independent Flow Modules:",
                            ("✓ flow-ch (ICHSIC) - Example implementation", 1),
                            ("✓ flow-sepa - Future SEPA flows", 1),
                            ("✓ flow-swift - Future SWIFT flows", 1),
                            ("✓ Each flow: Own configuration, XSD, XSLT", 1)
                        ])
    
    # Slide 12: Core Technologies
    create_two_column_slide(prs,
                           "Core Technologies & Infrastructure",
                           [
                               "Core Technologies:",
                               "• Java 21 (LTS)",
                               "• Spring Boot 3.4.1",
                               "• Apache Camel 4.1.0",
                               "• Maven 3.9+",
                               "• Docker",
                               "",
                               "Processing:",
                               "• Saxon 12.3 (XSLT)",
                               "• Hibernate 6.2.7",
                               "• Liquibase 4.25.1"
                           ],
                           [
                               "Infrastructure:",
                               "• IBM MQ (Enterprise JMS)",
                               "• Apache Kafka (Events)",
                               "• Redis (Cache)",
                               "• PostgreSQL (Database)",
                               "• Axway CFT (File Transfer)",
                               "• Hawtio (Monitoring)",
                               "",
                               "Standards:",
                               "• ISO 20022",
                               "• PACS.008 (Swiss)"
                           ])
    
    # Slide 13: Core Features
    create_section_slide(prs, "Core Features")
    
    # Slide 14: Message Processing Pipeline
    create_content_slide(prs,
                        "Message Processing Pipeline",
                        [
                            "Complete processing workflow:",
                            "",
                            "1. Input → Message reception from JMS queues",
                            "2. Archive → Save to NAS file system",
                            "3. Identify → Lookup referential data (Redis cache)",
                            "4. Validate → XSD schema validation (ISO 20022)",
                            "5. Transform → XSLT transformation",
                            "6. Log → Publish events to Kafka",
                            "7. Output → Final routing and delivery",
                            "",
                            "Each step includes error handling, metrics, and logging"
                        ])
    
    # Slide 15: Caching Strategy
    create_content_slide(prs,
                        "Caching Strategy",
                        [
                            "Spring caching for high performance:",
                            "",
                            "✓ Configurable TTL (default: 3600 seconds)",
                            "✓ Cache-aside pattern implementation",
                            "✓ Automatic cache invalidation",
                            "✓ Fallback to service calls on cache miss",
                            "✓ Sub-millisecond lookup times",
                            "✓ Supports Redis, Caffeine, EhCache backends",
                            "",
                            "Benefits:",
                            ("Reduced latency for referential lookups", 1),
                            ("Lower load on PostgreSQL database", 1),
                            ("Improved overall system throughput", 1)
                        ])
    
    # Slide 16: Payment Flow Processing
    create_section_slide(prs, "Multi-Flow Architecture")
    
    # Slide 17: ICHSIC Flow Example
    create_content_slide(prs,
                        "Example: ICHSIC Flow (One of Many)",
                        [
                            "Flow Code: ICHSIC (Switzerland IN SEPA SIC)",
                            "",
                            "Step 1: Message Reception",
                            ("IBM MQ → K-MQ-Starter → Generate flow ID", 1),
                            ("Archive to /opt/nas/CH/IN/{flowOccurId}/", 1),
                            "",
                            "Step 2: Identification & Enrichment",
                            ("K-Identification → Spring Cache → Referential Service", 1),
                            "",
                            "Step 3-5: Validate → Transform → Log → CFT",
                            "",
                            "Other flows follow the same pattern with different configs"
                        ])
    
    # Slide 18: Infrastructure & Deployment
    create_section_slide(prs, "Infrastructure & Deployment")
    
    # Slide 19: Docker Compose Stack
    create_content_slide(prs,
                        "Docker Compose Stack",
                        [
                            "Complete containerized infrastructure:",
                            "",
                            "• Zookeeper - Kafka coordination",
                            "• Kafka - Event streaming (ports 9092, 29092)",
                            "• PostgreSQL - Database (port 5432)",
                            "• Redis - Cache (port 6379)",
                            "• IBM MQ - Enterprise message broker (ports 1414, 9443)",
                            "• Axway CFT - File transfer (port 1761)",
                            "• Referential API - Config service (port 8099)",
                            "• Flow Apps - Multiple payment processors (scalable)",
                            "• Hawtio - Camel monitoring (port 8090)",
                            "",
                            "Network: pixel-v2-network (Bridge)"
                        ])
    
    # Slide 20: Apache Camel 4.1.0
    create_section_slide(prs, "Apache Camel & Spring Boot")
    
    # Slide 21: Apache Camel Features
    create_content_slide(prs,
                        "Apache Camel - Enterprise Integration",
                        [
                            "Why Apache Camel?",
                            ("300+ components for connectivity", 1),
                            ("Enterprise Integration Patterns (EIP)", 1),
                            ("Kamelet support for reusable templates", 1),
                            ("Cloud-native and lightweight", 1),
                            "",
                            "Key Features in PIXEL-V2:",
                            ("Content-Based Routing", 1),
                            ("Message Transformation (XSLT, JSON, XML)", 1),
                            ("Error Handling (dead letter, retry, circuit breaker)", 1),
                            ("Parallel Processing with thread pools", 1),
                            ("Idempotency & duplicate detection", 1)
                        ])
    
    # Slide 22: Spring Boot 3.4.1
    create_content_slide(prs,
                        "Spring Boot 3.4.1 - Modern Framework",
                        [
                            "Why Spring Boot?",
                            ("Auto-configuration for rapid development", 1),
                            ("Production-ready features (health, metrics)", 1),
                            ("Cloud-native support", 1),
                            ("Seamless Camel integration", 1),
                            "",
                            "Features in PIXEL-V2:",
                            ("Spring Boot Starter for Camel", 1),
                            ("Configuration management (profiles)", 1),
                            ("Spring Data JPA + PostgreSQL", 1),
                            ("Spring Data Redis for distributed caching", 1),
                            ("Spring Cache abstraction (Redis/Caffeine/EhCache)", 1),
                            ("Actuator endpoints (/health, /metrics)", 1)
                        ])
    
    # Slide 23: Hawtio Monitoring
    create_content_slide(prs,
                        "Hawtio - Real-Time Camel Monitoring",
                        [
                            "What is Hawtio?",
                            ("Lightweight web console for Java applications", 1),
                            ("First-class Apache Camel support", 1),
                            "",
                            "Key Features:",
                            ("Real-time route visualization", 1),
                            ("Message tracing and inspection", 1),
                            ("Performance metrics (throughput, latency)", 1),
                            ("JMX management", 1),
                            ("Route control (start/stop/suspend)", 1),
                            ("Thread pool monitoring", 1),
                            "",
                            "Access: http://localhost:8090/hawtio"
                        ])
    
    # Slide 24: Kaoto Visual Designer
    create_content_slide(prs,
                        "Kaoto - Visual Camel Route Designer",
                        [
                            "What is Kaoto?",
                            ("Low-code integration designer for Camel", 1),
                            ("Drag-and-drop route creation", 1),
                            "",
                            "Key Features:",
                            ("Visual design canvas", 1),
                            ("Browse 300+ Camel components", 1),
                            ("Kamelet designer", 1),
                            ("Auto-generate YAML/Java code", 1),
                            ("Template library for common patterns", 1),
                            "",
                            "Use Cases:",
                            ("Design new payment flows visually", 1),
                            ("Rapid prototyping", 1),
                            ("Developer onboarding", 1)
                        ])
    
    # Slide 25: CFT Integration
    create_content_slide(prs,
                        "Axway CFT - File Transfer Integration",
                        [
                            "Secure file exchange with partners",
                            "",
                            "Features:",
                            ("Automatic file routing and transformation", 1),
                            ("Multiple protocols (SFTP, HTTPS, PeSIT)", 1),
                            ("Acknowledgment and tracking", 1),
                            ("High availability and fault tolerance", 1),
                            "",
                            "Use Cases:",
                            ("Partner file exchange (banks, institutions)", 1),
                            ("Batch payment file transmission", 1),
                            ("Regulatory report submission", 1),
                            ("Archive and backup transfers", 1)
                        ])
    
    # Slide 26: Security & Compliance
    create_section_slide(prs, "Security & Compliance")
    
    # Slide 27: Security Features
    create_two_column_slide(prs,
                           "Security & Compliance",
                           [
                               "Security Features:",
                               "• Spring Security integration",
                               "• Password encryption",
                               "• Docker network isolation",
                               "• Service-level access control",
                               "• Redis authentication",
                               "• IBM MQ secure channels",
                               "• CFT encryption",
                               "",
                               "Data Protection:",
                               "• Secure PostgreSQL storage",
                               "• NAS file permissions",
                               "• Encrypted connections"
                           ],
                           [
                               "Compliance:",
                               "• ISO 20022 full compliance",
                               "• Complete Kafka audit trail",
                               "• NAS message archiving",
                               "• Regulatory retention",
                               "• Strict XSD validation",
                               "",
                               "Audit Capabilities:",
                               "• Event-driven logging",
                               "• Transaction tracking",
                               "• Performance metrics",
                               "• Error reporting"
                           ])
    
    # Slide 28: Project Statistics
    create_content_slide(prs,
                        "Project Statistics",
                        [
                            "Core Modules: 2 (framework, referential)",
                            "Flow Modules: Multiple (flow-ch/ICHSIC + others)",
                            "Custom Kamelets: 15+ reusable components",
                            "Infrastructure Components: 9 Docker services (with CFT)",
                            "",
                            "Technology Versions:",
                            ("Java 21 (Latest LTS)", 1),
                            ("Spring Boot 3.4.1", 1),
                            ("Apache Camel 4.1.0", 1),
                            ("PostgreSQL (Latest)", 1),
                            ("Redis (Latest)", 1),
                            ("Kafka 7.5.0", 1)
                        ])
    
    # Slide 29: Future Roadmap
    create_content_slide(prs,
                        "Future Roadmap",
                        [
                            "Short-term Enhancements:",
                            ("Add more payment flows (SEPA, SWIFT, etc.)", 1),
                            ("Performance: Optimize Spring caching strategies", 1),
                            ("Kaoto Integration: Visual route design", 1),
                            "",
                            "Medium-term Goals:",
                            ("Multi-country: Additional payment schemes", 1),
                            ("CFT Expansion: Enhanced file transfer scenarios", 1),
                            ("Resilience: Circuit breaker with Camel", 1),
                            "",
                            "Long-term Vision:",
                            ("AI/ML: Fraud detection integration", 1),
                            ("Real-time: Stream processing optimization", 1),
                            ("Camel Quarkus: Native compilation support", 1)
                        ])
    
    # Slide 30: Key Takeaways
    create_content_slide(prs,
                        "Key Takeaways",
                        [
                            "✓ Apache Camel 4.1.0: Enterprise integration patterns",
                            "✓ Spring Boot 3.4.1: Modern cloud-native framework",
                            "✓ Hawtio: Real-time Camel route monitoring",
                            "✓ Kaoto: Visual low-code integration designer",
                            "✓ IBM MQ + Kafka + CFT: Enterprise messaging",
                            "✓ Compliant: ISO 20022 standards adherence",
                            "✓ Maintainable: Modular Kamelet architecture",
                            "✓ Flexible: Configuration-driven flows"
                        ])
    
    # Slide 31: Thank You
    create_title_slide(prs,
                      "Thank You!",
                      "PIXEL-V2 - Enterprise Payment Processing, Simplified")
    
    # Save the presentation
    output_file = "PIXEL-V2-Presentation.pptx"
    prs.save(output_file)
    print(f"✓ PowerPoint presentation created successfully: {output_file}")
    print(f"  Total slides: {len(prs.slides)}")

if __name__ == "__main__":
    main()
